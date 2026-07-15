import hashlib
import os
import time
import urllib.parse
from io import BytesIO

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SLIDE_W = 10
SLIDE_H = 5.625
IMAGE_CACHE = ".ppt_image_cache"
# Hard limits: leave roughly 40 seconds for text/layout/save inside a 3-minute job.
MAX_TOTAL_IMAGE_SECONDS = 140
IMAGE_TIMEOUT_SECONDS = 12
AI_RETRIES = 1


def _image_prompts(slide_data):
    """
    Preferred JSON field:
      "image_prompts": ["enterprise digital strategy workshop", "cloud data dashboard"]

    The old image_keyword field remains supported. Separate alternatives with |.
    """
    raw = slide_data.get("image_prompts") or slide_data.get("image_keywords")
    if isinstance(raw, str):
        raw = [part.strip() for part in raw.split("|")]
    elif not isinstance(raw, list):
        raw = []

    title = str(slide_data.get("title", ""))
    content = str(slide_data.get("content", "")).replace("\n", " ")[:180]
    legacy = slide_data.get("image_keyword", "")
    if not raw and legacy:
        raw = [part.strip() for part in str(legacy).split("|")]

    prompts = []
    for item in raw:
        clean = str(item).strip()
        if clean and clean.casefold() not in [p.casefold() for p in prompts]:
            prompts.append(f"{title}；{clean}")

    # Chinese title/content are intentionally retained: the image model can use them.
    if not prompts:
        prompts.append(f"{title}；{content}")
    if len(prompts) == 1:
        prompts.append(f"{title}；{content}；以数据图形和业务场景为主的另一种构图")
    return prompts[:3]


def _valid_image(response):
    if response.status_code != 200 or len(response.content) < 4_000:
        return False
    content_type = response.headers.get("Content-Type", "").lower()
    if not content_type.startswith("image/"):
        return False
    # Common signatures: JPEG, PNG, GIF, WEBP.
    head = response.content[:12]
    return (
        head.startswith(b"\xff\xd8\xff")
        or head.startswith(b"\x89PNG\r\n\x1a\n")
        or head.startswith((b"GIF87a", b"GIF89a"))
        or (head.startswith(b"RIFF") and response.content[8:12] == b"WEBP")
    )


def _download(session, url, timeout=IMAGE_TIMEOUT_SECONDS):
    try:
        response = session.get(url, timeout=timeout, allow_redirects=True)
        return response.content if _valid_image(response) else None
    except requests.RequestException:
        return None


def _pollinations_url(prompt, seed, variant=0):
    styles = [
        "clean modern corporate editorial illustration",
        "premium isometric business technology illustration",
        "realistic corporate technology scene",
    ]
    detailed = (
        f"{prompt}, {styles[variant % len(styles)]}, one clear subject, "
        "wide composition, blue and orange accents, no text, no letters, no logo, no watermark"
    )
    encoded = urllib.parse.quote(detailed, safe="")
    return (
        f"https://image.pollinations.ai/prompt/{encoded}"
        f"?width=1000&height=650&nologo=true&seed={seed}"
    )


def fetch_relevant_image(session, prompt, slide_index, image_index, deadline=None):
    """Try the content-specific AI endpoint with different seeds/styles."""
    os.makedirs(IMAGE_CACHE, exist_ok=True)
    identity = f"v3|{slide_index}|{image_index}|{prompt}"
    cache_key = hashlib.sha256(identity.encode("utf-8")).hexdigest()[:20]
    cache_path = os.path.join(IMAGE_CACHE, cache_key + ".img")
    if os.path.exists(cache_path) and os.path.getsize(cache_path) > 4_000:
        with open(cache_path, "rb") as file:
            return BytesIO(file.read())

    if deadline is not None and time.monotonic() >= deadline:
        return None

    seed = int(hashlib.sha256(identity.encode("utf-8")).hexdigest()[:8], 16) % 999_999
    urls = [_pollinations_url(prompt, seed + retry, retry) for retry in range(AI_RETRIES)]
    for url in urls:
        remaining = deadline - time.monotonic() if deadline is not None else IMAGE_TIMEOUT_SECONDS
        if remaining <= 0.5:
            return None
        content = _download(session, url, timeout=min(IMAGE_TIMEOUT_SECONDS, remaining))
        if content:
            with open(cache_path, "wb") as file:
                file.write(content)
            return BytesIO(content)
        if AI_RETRIES > 1:
            time.sleep(0.2)
    return None


def _remove_layout_shapes(slide):
    for shape in list(slide.shapes):
        element = shape.element
        element.getparent().remove(element)


def _add_text(slide, text, x, y, w, h, size, color, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align
    return box


def _topic_nodes(title):
    title = str(title)
    if any(word in title for word in ("挑战", "风险", "问题")):
        return "关键挑战", ["安全", "人才", "治理"]
    if any(word in title for word in ("策略", "建议", "路径", "步骤")):
        return "行动路径", ["诊断", "实施", "迭代"]
    if any(word in title for word in ("未来", "趋势", "展望")):
        return "趋势演进", ["智能", "融合", "协作"]
    if any(word in title for word in ("技术", "核心")):
        return "技术体系", ["云端", "数据", "AI"]
    if any(word in title for word in ("案例", "成果", "总结")):
        return "价值闭环", ["技术", "业务", "增长"]
    return "业务图景", ["数据", "技术", "场景"]


def _add_topic_visual(slide, x, y, w, h, title, visual_index=0):
    """A relevant editable vector visual used only when all downloads fail."""
    heading, nodes = _topic_nodes(title)
    _add_text(
        slide, heading, x + 0.25, y + 0.18, w - 0.5, 0.28,
        10, RGBColor(71, 85, 105), True, PP_ALIGN.CENTER,
    )
    colors = [
        RGBColor(37, 99, 235),
        RGBColor(14, 165, 233),
        RGBColor(249, 115, 22),
    ]
    node_size = min(0.58, h * 0.34)
    gap = (w - 0.5 - node_size * 3) / 2
    start_x = x + 0.25
    node_y = y + max(0.62, (h - node_size) / 2 + 0.12)
    if visual_index % 2:
        nodes = list(reversed(nodes))
        colors = list(reversed(colors))
    for index, (node, color) in enumerate(zip(nodes, colors)):
        node_x = start_x + index * (node_size + gap)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(node_x), Inches(node_y),
            Inches(node_size), Inches(node_size),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        frame = circle.text_frame
        frame.clear()
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = frame.paragraphs[0]
        p.text = node
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(8 if node_size < 0.5 else 9)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        if index < 2:
            _add_text(
                slide, "›", node_x + node_size + 0.02, node_y + 0.08,
                max(0.16, gap - 0.04), 0.3, 16,
                RGBColor(148, 163, 184), True, PP_ALIGN.CENTER,
            )


def _add_image_card(slide, image_stream, x, y, w, h, title, visual_index=0):
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.04), Inches(y + 0.05), Inches(w), Inches(h),
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(220, 226, 235)
    shadow.line.fill.background()

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(245, 247, 250)
    card.line.color.rgb = RGBColor(224, 229, 237)

    if image_stream:
        try:
            # python-pptx keeps aspect ratio when only width is supplied.
            picture = slide.shapes.add_picture(
                image_stream, Inches(x + 0.04), Inches(y + 0.04), width=Inches(w - 0.08)
            )
            if picture.height > Inches(h - 0.08):
                max_height = Inches(h - 0.08)
                scale = max_height / picture.height
                picture.width = int(picture.width * scale)
                picture.height = max_height
            picture.left = Inches(x + (w - picture.width / Inches(1)) / 2)
            picture.top = Inches(y + (h - picture.height / Inches(1)) / 2)
            return
        except Exception:
            pass

    _add_topic_visual(slide, x, y, w, h, title, visual_index)


def _add_body(slide, content_text, x=0.55, y=1.35, w=5.2, h=3.85):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.04)
    lines = [line.strip() for line in str(content_text).splitlines() if line.strip()]
    for index, line in enumerate(lines or [""]):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(55, 65, 81)
        p.space_after = Pt(10)
        p.line_spacing = 1.15


def create_ppt_from_json(json_data, output_filename, template_path="template.pptx"):
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        while prs.slides:
            relationship_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(relationship_id)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    session = requests.Session()
    session.headers.update({"User-Agent": "PptImageBuilder/2.0 (presentation generator)"})
    image_deadline = time.monotonic() + MAX_TOTAL_IMAGE_SECONDS

    for slide_index, slide_data in enumerate(json_data.get("slides", [])):
        title = slide_data.get("title", "未命名页面")
        content = slide_data.get("content", "")
        prompts = _image_prompts(slide_data)

        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        _remove_layout_shapes(slide)

        # Accent line and title establish consistent hierarchy.
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.03), Inches(9), Inches(0.035)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(37, 99, 235)
        accent.line.fill.background()
        _add_text(slide, title, 0.5, 0.28, 9, 0.68, 25, RGBColor(31, 41, 55), True)

        if slide_index == 0:
            _add_text(
                slide, content, 0.7, 1.55, 4.25, 2.7, 17,
                RGBColor(75, 85, 99), align=PP_ALIGN.LEFT,
            )
            hero = fetch_relevant_image(
                session, prompts[0], slide_index, 0, image_deadline
            )
            _add_image_card(slide, hero, 5.25, 1.35, 4.15, 3.65, title, 0)
        else:
            _add_body(slide, content)
            # Two content-specific visuals per normal page.
            for image_index in range(2):
                prompt = prompts[image_index % len(prompts)]
                print(
                    f"第 {slide_index + 1} 页图片 {image_index + 1}/2：{prompt}",
                    flush=True,
                )
                stream = fetch_relevant_image(
                    session, prompt, slide_index, image_index, image_deadline
                )
                _add_image_card(
                    slide,
                    stream,
                    6.05,
                    1.3 + image_index * 2.02,
                    3.4,
                    1.82,
                    title,
                    image_index,
                )

    prs.save(output_filename)
    print(f"PPT 已保存：{output_filename}")


if __name__ == "__main__":
    # Replace this sample with your own JSON or import this module and call the function.
    sample = {
        "slides": [
            {
                "title": "数字化转型",
                "content": "以数据和技术重塑业务增长",
                "image_prompts": ["enterprise digital transformation strategy"],
            },
            {
                "title": "数字化转型的核心技术",
                "content": "1. 云计算：提供弹性基础设施\n2. 大数据分析：支持精准决策\n3. 人工智能：实现业务自动化",
                "image_prompts": [
                    "cloud computing data center",
                    "artificial intelligence business analytics dashboard",
                ],
            },
        ]
    }
    create_ppt_from_json(sample, "digital_transformation_demo.pptx")
